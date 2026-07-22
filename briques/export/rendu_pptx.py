"""Rendu PPTX DÉTERMINISTE (python-pptx) — diapositives structurées → fichier PPTX.

python-pptx est pur-Python (pas de libs système type Pango/Cairo requises comme pour
WeasyPrint, cf. rendu_pdf.py) : ce module peut être appelé pour de vrai dans les tests
offline, sans mock.

Thème v1 unique : 'sobre' (texte seul, layout titre+contenu standard). Pas de génération
d'images/graphiques (hors périmètre v1, cf. design S194).
"""
from __future__ import annotations

import io

from pptx import Presentation

THEMES = {"sobre"}


def generer(titre: str, diapositives: list[dict], theme: str = "sobre") -> bytes:
    if not (titre or "").strip():
        raise ValueError("Le titre est vide.")
    if not diapositives:
        raise ValueError("Aucune diapositive fournie.")
    if theme not in THEMES:
        raise ValueError(f"Thème inconnu : {theme!r} (attendu : {sorted(THEMES)}).")

    prez = Presentation()
    diapo_titre = prez.slides.add_slide(prez.slide_layouts[0])
    diapo_titre.shapes.title.text = titre

    layout_contenu = prez.slide_layouts[1]   # layout standard "Title and Content"
    for d in diapositives:
        diapo = prez.slides.add_slide(layout_contenu)
        diapo.shapes.title.text = (d.get("titre") or "").strip() or "—"
        points = d.get("points") or []
        corps = diapo.placeholders[1].text_frame
        if points:
            corps.text = points[0]
            for point in points[1:]:
                corps.add_paragraph().text = point
        notes = (d.get("notes") or "").strip()
        if notes:
            diapo.notes_slide.notes_text_frame.text = notes

    buffer = io.BytesIO()
    prez.save(buffer)
    return buffer.getvalue()
