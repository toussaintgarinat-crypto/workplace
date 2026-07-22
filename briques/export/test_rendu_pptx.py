"""Tests — rendu_pptx.py (python-pptx est pur-Python : ces tests appellent le VRAI rendu,
pas de mock — cf. rendu_pdf.py pour la raison du mock côté WeasyPrint)."""
import io

from pptx import Presentation

import rendu_pptx as PP


def test_generer_titre_vide_leve():
    try:
        PP.generer("", [{"titre": "A", "points": ["x"]}])
    except ValueError as e:
        assert "titre" in str(e).lower()
    else:
        raise AssertionError("titre vide aurait dû lever ValueError")


def test_generer_sans_diapositives_leve():
    try:
        PP.generer("Deck", [])
    except ValueError as e:
        assert "diapositive" in str(e).lower()
    else:
        raise AssertionError("diapositives vides auraient dû lever ValueError")


def test_generer_theme_inconnu_leve():
    try:
        PP.generer("Deck", [{"titre": "A", "points": ["x"]}], theme="neon")
    except ValueError as e:
        assert "thème" in str(e).lower() or "theme" in str(e).lower()
    else:
        raise AssertionError("thème inconnu aurait dû lever ValueError")


def test_generer_produit_un_pptx_valide_relisible():
    data = PP.generer("Mon Deck", [
        {"titre": "Diapo 1", "points": ["Point A", "Point B"], "notes": "note interne"},
        {"titre": "Diapo 2", "points": ["Point C"]},
    ])
    assert data[:2] == b"PK"   # signature ZIP/OOXML : c'est un vrai fichier, pas un mock
    prez = Presentation(io.BytesIO(data))
    diapos = list(prez.slides)
    assert len(diapos) == 3   # 1 diapo de titre + 2 diapos de contenu
    assert diapos[0].shapes.title.text == "Mon Deck"
    assert diapos[1].shapes.title.text == "Diapo 1"
    corps = diapos[1].placeholders[1].text_frame
    textes = [p.text for p in corps.paragraphs]
    assert textes == ["Point A", "Point B"]
    assert diapos[1].notes_slide.notes_text_frame.text == "note interne"
    assert diapos[2].shapes.title.text == "Diapo 2"
